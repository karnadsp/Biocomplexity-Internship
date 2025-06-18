from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pathlib import Path
import json
import os
import re

def auto_size_text(text_frame, min_size=8, max_size=24):
    """Auto-size text to fit within the text frame"""
    # Enable auto-size
    text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    
    # Get the paragraph
    paragraph = text_frame.paragraphs[0]
    
    # Set initial font size
    paragraph.font.size = Pt(max_size)
    
    # If text is too long, reduce font size
    while len(text_frame.text) > 1000 and paragraph.font.size > Pt(min_size):
        current_size = paragraph.font.size.pt
        paragraph.font.size = Pt(current_size - 1)

def format_text_for_slide(text, max_length=1000):
    """Format text to be more suitable for slides"""
    # If text is too long, truncate and add ellipsis
    if len(text) > max_length:
        text = text[:max_length] + "..."
    
    # Replace multiple newlines with single newline
    text = re.sub(r'\n\s*\n', '\n\n', text)
    
    # Ensure text ends with a single newline
    text = text.strip() + '\n'
    
    return text

def create_title_slide(prs, experiment_name, timestamp):
    """Create the title slide"""
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = f"Biological System Modeling Experiment"
    subtitle.text = f"{experiment_name}\n{timestamp}"
    
    # Auto-size the title
    if title.text_frame:
        auto_size_text(title.text_frame, min_size=24, max_size=44)
    
    # Auto-size the subtitle
    if subtitle.text_frame:
        auto_size_text(subtitle.text_frame, min_size=18, max_size=32)

def create_description_slide(prs, description_data):
    """Create a slide with the initial description"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "System Description"
    content.text = format_text_for_slide(description_data['output']['description'])
    
    # Auto-size the content
    if content.text_frame:
        auto_size_text(content.text_frame)

def create_ontology_slide(prs, ontology_data):
    """Create a slide with ontology annotations"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Ontology Annotations"
    
    # Format the annotations
    text_frame = content.text_frame
    text_frame.clear()
    
    response = ontology_data['output']['response']
    if '```json' in response:
        response = response.split('```json')[1].split('```')[0].strip()
    
    try:
        annotations = json.loads(response)
        formatted_text = []
        
        for category, items in annotations.items():
            if items:  # Only show non-empty categories
                formatted_text.append(f"\n{category}:")
                
                if isinstance(items, list):
                    for item in items:
                        if 'CellOntology' in item:
                            formatted_text.append(f"• {item['CellOntology']['label']} (ID: {item['CellOntology']['id']})")
                        elif 'id' in item:
                            formatted_text.append(f"• {item['label']} (ID: {item['id']})")
                elif isinstance(items, dict):
                    for subcat, subitems in items.items():
                        if subitems:
                            formatted_text.append(f"  {subcat}:")
                            for subitem in subitems:
                                formatted_text.append(f"  • {subitem}")
        
        content.text = format_text_for_slide('\n'.join(formatted_text))
        
    except json.JSONDecodeError:
        print(f"Warning: Could not parse JSON from response: {response}")
        content.text = format_text_for_slide(response)
    
    # Auto-size the content
    if content.text_frame:
        auto_size_text(content.text_frame)

def create_clarifications_slide(prs, clarifications_data):
    """Create a slide with clarifications"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Clarifications"
    content.text = format_text_for_slide(clarifications_data['output']['answers'])
    
    # Auto-size the content
    if content.text_frame:
        auto_size_text(content.text_frame)

def create_llm_response_slide(prs, response_data):
    """Create a slide for general LLM responses"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    # Extract step name from the data or use a default
    step_name = response_data.get('step', 'LLM Response').replace('_', ' ').title()
    title.text = step_name
    
    # Format the response - handle different response formats
    output = response_data.get('output', {})
    if 'response' in output:
        content.text = format_text_for_slide(output['response'])
    elif 'answers' in output:
        content.text = format_text_for_slide(output['answers'])
    elif 'description' in output:
        content.text = format_text_for_slide(output['description'])
    else:
        # If no known format is found, display the entire output as JSON
        content.text = format_text_for_slide(json.dumps(output, indent=2))
    
    # Auto-size the content
    if content.text_frame:
        auto_size_text(content.text_frame)

def create_presentation(experiment_dir):
    """Create a PowerPoint presentation from individual JSON files in the experiment directory"""
    # Create presentation
    prs = Presentation()
    
    # Set slide dimensions to 16:9
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # Get experiment name and timestamp from directory name
    experiment_name = experiment_dir.name
    timestamp = experiment_name.split('_')[-1] if '_' in experiment_name else ""
    
    # Create title slide
    create_title_slide(prs, experiment_name, timestamp)
    
    # Process all JSON files in the experiment directory
    json_files = sorted(experiment_dir.glob("*.json"))
    
    for json_file in json_files:
        try:
            with open(json_file, 'r') as f:
                data = json.load(f)
            
            # Skip if data is empty or not a dictionary
            if not data or not isinstance(data, dict):
                print(f"Warning: Invalid data in {json_file}")
                continue
            
            # Determine the type of slide to create based on the step name
            step = data.get('step', '')
            
            if step == 'initial_description':
                create_description_slide(prs, data)
            elif step == 'clarifications':
                create_clarifications_slide(prs, data)
            elif 'ontology' in step.lower():
                create_ontology_slide(prs, data)
            else:
                create_llm_response_slide(prs, data)
                
        except json.JSONDecodeError:
            print(f"Warning: Could not parse JSON file: {json_file}")
            continue
        except Exception as e:
            print(f"Warning: Error processing {json_file}: {str(e)}")
            continue
    
    # Save the presentation
    output_file = Path(experiment_dir) / "experiment_presentation.pptx"
    prs.save(output_file)
    return output_file

def main():
    # Get the most recent experiment directory
    experiments_dir = Path("experiments")
    if not experiments_dir.exists():
        print("No experiments directory found!")
        return
    
    # Get the most recent experiment
    experiment_dirs = sorted(experiments_dir.glob("*"), key=lambda x: x.stat().st_mtime, reverse=True)
    if not experiment_dirs:
        print("No experiments found!")
        return
    
    latest_experiment = experiment_dirs[-0]
    print(f"Creating presentation for experiment: {latest_experiment.name}")
    
    # Create the presentation
    output_file = create_presentation(latest_experiment)
    print(f"Presentation created successfully: {output_file}")

if __name__ == "__main__":
    main() 